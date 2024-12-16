import os
import shutil
import random
from loremipsum import get_paragraphs
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
import zipfile

# vars
FILE_EXTENSIONS = ('docx', 'xlsx', 'txt', 'config', 'xml', 'pdf', 'zip')
BASE_DIR = './FILES'
FILES_PER_DIR = 100
MIN_PARAGRAPHS = 1
MAX_PARAGRAPHS = 10

# Corporate bs
CORPORATE_WORDS = [
    "Synergy", "Innovate", "Optimize", "Strategize", "Transform", 
    "Paradigm", "Leverage", "Enterprise", "Vision", "Empower", 
    "Matrix", "Dynamic", "Core", "Solution", "Global", 
    "Network", "Growth", "Efficiency", "Performance", "Insight",
    "Analytics", "Disrupt", "Agile", "Scalability", "Engagement",
    "Innovation", "Value", "Alignment", "Bandwidth", "Ecosystem",
    "Vertical", "Horizontal", "Stakeholder", "KPI", "ROI", 
    "Benchmark", "BestPractices", "ChangeManagement", "Collaborate",
    "Crossfunctional", "Deliverable", "Empowerment", "Feedback",
    "Gamification", "Holistic", "Ideation", "Journey", "Key",
    "Leadership", "Mindset", "Navigate", "Operational", "Pivot"
]

# Load word list
with open('words_alpha.txt', 'r') as f:
    WORD_LIST = [line.strip() for line in f]

def get_random_word(word_list):
    return random.choice(word_list)

def generate_file_content(paragraph_count):
    paragraphs = get_paragraphs(random.randint(MIN_PARAGRAPHS, paragraph_count))
    return '\n\n'.join(paragraphs).replace("b'", "").replace("B'", "").replace("'", "")

def select_file_extension(extensions):
    return random.choice(extensions)

def generate_file_name(extension):
    return f"{get_random_word(WORD_LIST)}_{get_random_word(WORD_LIST)}{random.randint(1, 9)}.{extension}"

def generate_directory_name():
    # Randomly choose 2 to 4 corpo bs words and mash them together.
    num_words = random.randint(2, 4)
    return '_'.join(random.sample(CORPORATE_WORDS, num_words))

def create_directory_structure(base_path):
    print("Creating Directory Tree...")
    shutil.rmtree(base_path, ignore_errors=True)
    os.makedirs(base_path)
    # 26 directories with corporate buzzwords
    for _ in range(26):
        dir_name = generate_directory_name()
        os.makedirs(os.path.join(base_path, dir_name))
        print(f"Created Directory: {dir_name}")

def write_file_content(path, content, extension):
    if extension in ['txt', 'config']:
        with open(path, 'w') as f:
            f.write(content)
    elif extension == 'docx':
        doc = Document()
        doc.add_paragraph(content)
        doc.save(path)
    elif extension == 'xlsx':
        wb = Workbook()
        ws = wb.active
        for line in content.split('\n'):
            ws.append([line])
        wb.save(path)
    elif extension == 'xml':
        with open(path, 'w') as f:
            f.write('<?xml version="1.0" encoding="UTF-8"?>\n<content>\n' + content + '\n</content>')
    elif extension == 'pptx':
        prs = Presentation()
        for line in content.split('\n')[:100]:  # Limit to 100 slides for performance
            if line:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(5))
                tf = txBox.text_frame
                tf.text = line
        prs.save(path)
    elif extension == 'pdf':
        p = canvas.Canvas(path, pagesize=letter)
        width, height = letter
        for line in content.split('\n'):
            p.drawString(1*inch, height - (1 + content.split('\n').index(line))*inch, line)
        p.save()
    elif extension == 'zip':
        with zipfile.ZipFile(path, 'w') as zipf:
            for i in range(3):
                dummy_name = generate_file_name('txt')
                dummy_path = os.path.join(BASE_DIR, 'temp', dummy_name)
                os.makedirs(os.path.dirname(dummy_path), exist_ok=True)
                with open(dummy_path, 'w') as dummy_file:
                    dummy_file.write(generate_file_content(1))
                zipf.write(dummy_path, dummy_name)
                os.remove(dummy_path)
            shutil.rmtree(os.path.join(BASE_DIR, 'temp'), ignore_errors=True)  # Remove temp directory

def main():
    print("-------------------")
    print("Generate Crap v1.1")
    print("Created by Andy Swift")
    print("Twitter: @SwiftSecur1")
    print("GitHub: https://github.com/SwiftSecur")
    print("-------------------")
    create_directory_structure(BASE_DIR)
    print(f"Generating {FILES_PER_DIR} files per directory...")

    for dir_name in os.listdir(BASE_DIR):
        dir_path = os.path.join(BASE_DIR, dir_name)
        if os.path.isdir(dir_path):
            print(f"Generating Files in Directory: {dir_name}")
            for _ in range(FILES_PER_DIR):
                extension = select_file_extension(FILE_EXTENSIONS)
                file_name = generate_file_name(extension)
                paragraph_count = random.randint(MIN_PARAGRAPHS, MAX_PARAGRAPHS)
                content = generate_file_content(paragraph_count)

                file_path = os.path.join(dir_path, file_name)
                write_file_content(file_path, content, extension)
                print(f"\t{file_name}\t{paragraph_count} paragraphs")

if __name__ == "__main__":
    main()
